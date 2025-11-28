from fastapi import FastAPI, HTTPException
from fastapi.responses import HTMLResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from contextlib import asynccontextmanager
from playwright.async_api import async_playwright, Page
import httpx
import asyncio
from dotenv import load_dotenv
from os import getenv
from pathlib import Path
import re
import pytz
from datetime import datetime
from urllib.parse import urlparse, urljoin
import json
from typing import Any
import uvicorn
from io import BytesIO
import pdfplumber
from docx import Document
from pptx import Presentation
import pandas as pd
import zipfile
import tarfile
import gzip
import base64

IST = pytz.timezone('Asia/Kolkata')

@asynccontextmanager
async def lifespan(app: FastAPI):
	app.state.playwright = await async_playwright().start()
	app.state.browser = await app.state.playwright.chromium.launch(headless=True)
	app.state.page = await app.state.browser.new_page()
	print("--- Browser launched ---")
	try:
		yield
	finally:
		await app.state.page.close()
		await app.state.browser.close()
		await app.state.playwright.stop()
		print("--- Browser closed ---")

app = FastAPI(lifespan=lifespan)

app.add_middleware(
	CORSMiddleware,
	allow_origins=["*"],
	allow_credentials=True,
	allow_methods=["*"],
	allow_headers=["*"],
)

# Configuration
templates_dir = Path(__file__).parent / "templates"
app.mount("/templates", StaticFiles(directory=str(templates_dir)), name="templates")
load_dotenv()
app.state.SECRET = getenv("SECRET")
app.state.LLM_API_KEY = getenv("LLM_API_KEY")
app.state.EMAIL = getenv("EMAIL")
app.state.LLM_BASE_URL = getenv("LLM_BASE_URL")
app.state.LLM_MODEL = getenv("LLM_MODEL")

@app.get("/", response_class=HTMLResponse)
def read_root():
    """Serve the index.html page"""
    with open(templates_dir / "index.html") as f:
        return HTMLResponse(content=f.read())

def validate_secret(secret: str) -> bool:
    return secret == app.state.SECRET

def validate_email(email: str) -> bool:
    return email == app.state.EMAIL

def validate_url(url: str) -> bool:
    ''' validate url return True or string with error message '''
    try:
        result = urlparse(url)
        if result.scheme != "https" or not result.netloc:
            return "Invalid URL format"
        return True
    except Exception:
        return "Invalid URL format"
    
def llm_headers() -> dict:
    return {
        "Authorization": f"Bearer {app.state.LLM_API_KEY}",
        "Content-Type": "application/json"
    }

def clean_json_text(raw: str) -> str:
	"""Clean malformed JSON in <pre> blocks."""
	# Remove HTML tags like <span class="origin">...</span>
	raw = re.sub(r"<[^>]+>", "", raw)

	# Replace invalid ellipsis (...) with null
	raw = raw.replace("...", "null")

	# Remove trailing commas before closing braces
	raw = re.sub(r",\s*([}\]])", r"\1", raw)

	return raw.strip()

async def extract_everything(page: Page, url: str):
    """Load a quiz URL and extract all data for LLM."""
    # load the page
    await page.goto(url, wait_until="networkidle")
    # extract full HTML
    html = await page.content()
    print(f"--- Extracted HTML content from {url} ---")

    page_screenshot = await page.screenshot()
    print(f"--- Captured screenshot of {url} ---")
    # extract text from screenshot using llm_api_key and gemini-2.0-flash
    # send to llm for image description
    system_prompt = "You are an automated image description agent. Describe the content of the provided image accurately."
    user_prompt = "Describe and extract content of the following image."
    payload = {
        "model": "gemini-2.0-flash",
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user",
             "content": [
                 {"type": "text", "text": user_prompt},
                 {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64.b64encode(page_screenshot).decode()}"}}
             ]}
        ],
        "temperature": 0.1
    }
    async with httpx.AsyncClient(timeout=30) as client:
        response = await client.post(
            f"{app.state.LLM_BASE_URL}/chat/completions", 
            json=payload, 
            headers=llm_headers(),
            timeout=30
        )
        response.raise_for_status()
    result = response.json()
    screenshot_description = result['choices'][0]['message']['content'].strip()
    print(f"--- Extracted Screenshot Description from {url} ---")

    # extract pre, code blocks as JSON payload templates
    precode = []
    blocks = await page.query_selector_all("pre, code")

    for block in blocks:
        raw = (await block.inner_text()).strip()

        # try raw JSON
        try:
            precode.append(json.loads(raw))
            continue
        except:
            pass

        # clean JSON and retry
        cleaned = clean_json_text(raw)
        try:
            precode.append(json.loads(cleaned))
        except:
            pass
    print(f"--- Extracted {len(precode)} pre/code blocks from {url} ---")

    # url extract, remove query params and fragments
    parsed_url = urlparse(page.url)
    base_url = f"{parsed_url.scheme}://{parsed_url.netloc}{parsed_url.path}"

    # extract scr, href, links from page
    links = []
    a_tags = await page.query_selector_all("a")
    audio_tags = await page.query_selector_all("audio")
    img_tags = await page.query_selector_all("img")
    video_tags = await page.query_selector_all("video")
    script_tags = await page.query_selector_all("script")
    link_tags = await page.query_selector_all("link")
    for a in a_tags:
        href = await a.get_attribute("href")
        if href:
            links.append(urljoin(base_url, href))
    for audio in audio_tags:
        src = await audio.get_attribute("src")
        if src:
            links.append(urljoin(base_url, src))
    for img in img_tags:
        src = await img.get_attribute("src")
        if src:
            links.append(urljoin(base_url, src))
    for video in video_tags:
        src = await video.get_attribute("src")
        if src:
            links.append(urljoin(base_url, src))
    for script in script_tags:
        src = await script.get_attribute("src")
        if src:
            links.append(urljoin(base_url, src))
    for link in link_tags:
        href = await link.get_attribute("href")
        if href:
            links.append(urljoin(base_url, href))
    print(f"--- Extracted {len(links)} links from {url} ---")

    # populate liks to appropriate categories using extensions
    pdfs, audios, images, linked_pages, others = [], [], [], {}, []
    for h in links:
        if h.endswith(".pdf"):
            # use pdfplumber to extract text, tables from pdf using link
            try:
                async with httpx.AsyncClient(timeout=30) as client:
                    response = await client.get(h)
                response.raise_for_status()
                with pdfplumber.open(BytesIO(response.content)) as pdf:
                    text = ""
                    tables = {}
                    for i, page in enumerate(pdf.pages, start=1):
                        text += f"Page {i}" + page.extract_text() + "\n"
                        for table in page.extract_tables():
                            tables[f"Page_{i}_Table_{len(tables)+1}"] = table
                pdfs.append({
                    "url": h,
                    "text": text.strip(),
                    "tables": tables
                })
                print(f"--- Extracted PDF content from {h} ---")
            except Exception as e:
                print(f"--- PDF Extraction Error for {h}: {e} ---")
            
        elif any(h.lower().endswith(ext) for ext in [".mp3", ".wav", ".flac", ".aac", ".ogg", ".m4a", ".alac", ".wma"]):
            # use llm_api_key to extract audio transcription from link using gpt-4o-transcribe
            # use link to load audio file and convert to bytes using BytesIO
            try:
                async with httpx.AsyncClient(timeout=30) as client:
                    response = await client.get(h)
                response.raise_for_status()
                audio_bytes = BytesIO(response.content)
                # send to llm for transcription
                system_prompt = "You are an automated transcription agent. Transcribe the provided audio file into text accurately."
                user_prompt = "Transcribe the following audio file into text."
                payload = {
                    "model": "gpt-4o-transcribe",
                    "messages": [
                        {"role": "system", "content": system_prompt},
                        {"role": "user",
                         "content": [
                            {"type": "text", "text": user_prompt},
                            {"type": "audio_url", "audio_url": {"url": f"data:audio/{h.split('.')[-1]};base64," + base64.b64encode(audio_bytes.getvalue()).decode()}}
                         ]}
                    ],
                    "temperature": 0.1
                }
                async with httpx.AsyncClient(timeout=30) as client:
                    response = await client.post(
                        f"{app.state.LLM_BASE_URL}/chat/completions", 
                        json=payload, 
                        headers=llm_headers(),
                        timeout=30
                    )
                response.raise_for_status()
                result = response.json()
                transcription = result['choices'][0]['message']['content'].strip()
                print(f"--- Extracted Audio Transcription from {h} ---")
                audios.append({
                    "url": h,
                    "transcription": transcription
                })
            except Exception as e:
                print(f"--- Audio Transcription Error for {h}: {e} ---")
                transcription = "N/A"
            
        elif any(h.lower().endswith(ext) for ext in [".jpg", ".jpeg", ".png", ".gif", ".webp", ".svg", ".avif"]):
            # use llm_api_key to extract image description from link using gemini-2.0-flash
            # use link to load image file and convert to bytes using BytesIO
            try:
                async with httpx.AsyncClient(timeout=30) as client:
                    response = await client.get(h)
                response.raise_for_status()
                image_bytes = BytesIO(response.content)
                # send to llm for image description
                system_prompt = "You are an automated image description agent. Describe the content of the provided image accurately."
                user_prompt = "Describe and extract content of the following image."
                payload = {
                    "model": "gemini-2.0-flash",
                    "messages": [
                        {"role": "system", "content": system_prompt},
                        {"role": "user",
                         "content": [
                            {"type": "text", "text": user_prompt},
                            {"type": "image_url", "image_url": {"url": f"data:image/{h.split('.')[-1]};base64," + base64.b64encode(image_bytes.getvalue()).decode()}}
                         ]},
                    ],
                    "temperature": 0.1
                }
                async with httpx.AsyncClient(timeout=30) as client:
                    response = await client.post(
                        f"{app.state.LLM_BASE_URL}/chat/completions", 
                        json=payload, 
                        headers=llm_headers(),
                        timeout=30
                    )
                response.raise_for_status()
                result = response.json()
                description = result['choices'][0]['message']['content'].strip()
                print(f"--- Extracted Image Description from {h} ---")
                images.append({
                    "url": h,
                    "description": description
                })
            except Exception as e:
                print(f"--- Image Description Error for {h}: {e} ---")
                description = "N/A"
            
        elif h.endswith(".html") or h.endswith(".htm") or re.match(r'.*/[^./]+$', h):
            # treat as linked page
            await page.goto(h, wait_until="networkidle")
            l_html = await page.content()
            linked_pages[h] = {
                "html": l_html
            }
            print(f"--- Extracted Linked Page HTML content from {h} ---")
            # restore original page
            await page.goto(url, wait_until="networkidle")
        else:
            if h.endswith(('.css', '.js', '.json', '.txt')):
                others.append({
                    "url": h,
                    "type": "file"
                })
            elif h.endswith(('.zip', '.tar', '.tar.gz', '.tgz', '.gz')):
                # extract all files from archive files
                try:
                    async with httpx.AsyncClient(timeout=30) as client:
                        response = await client.get(h)
                    response.raise_for_status()
                    file_bytes = BytesIO(response.content)
                    extracted_files = {}
                    if h.endswith('.zip'):
                        with zipfile.ZipFile(file_bytes) as z:
                            for file_info in z.infolist():
                                with z.open(file_info) as f:
                                    extracted_files[file_info.filename] = f.read().decode(errors='ignore')
                    elif h.endswith(('.tar', '.tar.gz', '.tgz')):
                        mode = 'r:gz' if h.endswith(('.tar.gz', '.tgz')) else 'r'
                        with tarfile.open(fileobj=file_bytes, mode=mode) as t:
                            for member in t.getmembers():
                                f = t.extractfile(member)
                                if f:
                                    extracted_files[member.name] = f.read().decode(errors='ignore')
                    elif h.endswith('.gz'):
                        with gzip.open(file_bytes, 'rt', errors='ignore') as f:
                            for line in f:
                                extracted_files[f.name] = line
                    print(f"--- Extracted Archive content from {h} ---")
                    others.append({
                        "url": h,
                        "type": "archive",
                        "files": extracted_files
                    })
                except Exception as e:
                    print(f"--- Archive Extraction Error for {h}: {e} ---")

            elif h.endswith(('.docx', '.xlsx', '.pptx', '.xls', '.doc')):
                # extract text from document files
                try:
                    async with httpx.AsyncClient(timeout=30) as client:
                        response = await client.get(h)
                    response.raise_for_status()
                    file_bytes = BytesIO(response.content)
                    if h.endswith(('.docx', '.doc')):
                        doc = Document(file_bytes)
                        full_text = []
                        for para in doc.paragraphs:
                            full_text.append(para.text)
                        text = '\n'.join(full_text)
                        others.append({
                            "url": h,
                            "type": "document",
                            "text": text.strip()
                        })
                    elif h.endswith('.pptx'):
                        prs = Presentation(file_bytes)
                        full_text = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    full_text.append(shape.text)
                        text = '\n'.join(full_text)
                        others.append({
                            "url": h,
                            "type": "presentation",
                            "text": text.strip()
                        })
                    elif h.endswith(('.xlsx', '.xls')):
                        df = pd.read_excel(file_bytes, engine='openpyxl')
                        df_csv = df.to_csv(index=False)
                        print(f"--- Extracted Document content from {h} ---")
                        others.append({
                            "url": h,
                            "type": "csv_converted",
                            "csv": df_csv
                        })
                    elif h.endswith('.csv'):
                        df = pd.read_csv(file_bytes)
                        df_csv = df.to_csv(index=False)
                        print(f"--- Extracted Document content from {h} ---")
                        others.append({
                            "url": h,
                            "type": "csv",
                            "csv": df_csv
                        })
                except Exception as e:
                    print(f"--- Document Extraction Error for {h}: {e} ---")

            elif h.endswith(('.mp4', '.mkv', '.mov', '.avi', '.wmv', '.flv', '.webm', '.mpeg')):
                others.append({
                    "url": h,
                    "type": "Video"
                })

    return {
        "html": html,
        "page_screenshot_data": screenshot_description,
        "pre_code_blocks": precode,
        "pdfs": pdfs,
        "audios": audios,
        "images": images,
        "linked_pages": linked_pages,
        "other_files": others
    }

async def solve_with_llm(content: dict[Any], previous_reason: str = None) -> str:
    """
    Sends scraped content to LLM to extract the question and find the answer.
    """
    system_prompt = (
        "You are an automated agent solving a quiz as a human. "
        "Your goal is to understand the question from the provided structured content "
        "and provide the specific answer. "
        "Return ONLY the answer. If it is a number, return just the digits. "
        "Do not provide email, secret, url, ignore them, provide only the answer field. "
        "If it is a number, return just the number. "
        "If it is a word, return just the word. "
        "If it is a phrase, return just the phrase without any quotes. "
        "If it is data, return JSON format with keys and values (values maybe data URI/attachment needed for 'answer''). "
        "If python script is needed, return only the script without any explanation that gives answer upon execution."
    )
    # pass dict as string
    user_prompt = f"Structured Content:\n{json.dumps(content, indent=2)}\n"
    
    if previous_reason:
        user_prompt += f"\n\nPREVIOUS ATTEMPT FAILED. Reason: {previous_reason}. Try a different approach."

    payload = {
        "model": app.state.LLM_MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        "temperature": 0.1
    }
    try:
        async with httpx.AsyncClient(timeout=150) as client:
            response = await client.post(
                f"{app.state.LLM_BASE_URL}/chat/completions", 
                json=payload, 
                headers=llm_headers(),
                timeout=150
            )
        response.raise_for_status()
        result = response.json()
        answer = result['choices'][0]['message']['content'].strip()
        answer = answer.replace("`", "").replace("'", "").replace('"', "")
        return answer
    except Exception as e:
        print(f"--- LLM Error: {e} ---")
        return None

async def process_task(url: str, reason: str = None) -> dict:
    """
    1. Scrape the data from URL or nested URLs
    2. Extract relevant content from various file types
    (Usually headless browser + file parsers)
    (html, script, css files of url and nested urls,
    audio(.mp3, .wav, .flac, .aac, .ogg, .m4a, .alac, and .wma),
    image(.jpg, .jpeg, .png, .gif, .webp, .svg, .avif),
    video(.mkv, .mp4, .mov, .avi, .wmv, .flv, .webm, .mpeg),
    csv, xlxs, docx, text, pdf(text, ocr, tables), zip files, tar files etc.)
    3. Understand the question/task
    4. Process it
    5. Return the answer
    """

    print(f"--- [{datetime.now(IST)}] Scraping: {url} ---")

    # 1. Scrape
    scraped_content = await extract_everything(app.state.page, url)
    if not scraped_content:
        print(f"--- Failed to scrape content from {url} ---")
        return None

    # 2. Solve with LLM
    print(f"--- Analyzing with LLM... ---")
    answer = await solve_with_llm(scraped_content, reason)
    print(f"--- Calculated Answer: {answer} ---")

    return {
            "email": app.state.EMAIL,
            "secret": app.state.SECRET,
            "url": url,
            "answer": answer
        }

async def worker(url: str):
    """
    Worker function to handle the task processing and submission loop.
    """

    max_iterations = 400
    iteration_count = 0
    questions_data = []
    questions_data.append({
        "url": url,
        "question_number": 1 + len(questions_data),
        "status": False,
        "started_at": datetime.now(IST),
        "completed_at": None,
        "answer": None
    })
    print("=== START EVALUATION ===\n")
    domain_match = re.search(r'^(?:http|ftp)s?://([^/]+)', url)
    domain = domain_match.group(1) if domain_match else 'default'

    # initial processing
    payload = await process_task(url)
    while True:
        # subsequent processing based on response  
        if not payload:
            print(f"--- Failed to generate payload for submission. Ending evaluation. ---\n")
            print("=== END EVALUATION DUE TO PAYLOAD ERROR ===\n")
            return None        
        try:
            async with httpx.AsyncClient(timeout=30) as client:
                response = await client.post(f"https://{domain}/submit", json=payload)
        except Exception as e:
            print(f"--- NETWORK ERROR ---\nTYPE: {type(e)}\nDETAILS: {str(e)}\n--- END NETWORK ERROR ---\n")
            print(f"=== END EVALUATION DUE TO NETWORK ERROR ===")
            return None
        # error, server error status codes handling
        if response.status_code in [400, 403, 404, 500, 502, 503, 504]:
            print(f"--- ERROR ---\nTYPE: HTTP {response.status_code}\nDETAILS: {response.text}\n--- END ERROR ---\n")
            print(f"=== END EVALUATION DUE TO ERROR ===")
            return None
        
        # handle is response not json convertable
        try:
            response.json()
        except Exception as e:
            print(f"--- ERROR ---\nTYPE: Invalid JSON Response\nDETAILS: {response.text}\n--- END ERROR ---\n")
            print(f"=== END EVALUATION DUE TO INVALID JSON RESPONSE ===")
            return None
        data = response.json()
        last_question = questions_data[-1]
        questions_data[-1]["answer"] = payload.get("answer", None)

        print(f"--- Answer submitted successfully to {domain}/submit ---\n")
        print(f"--- Submission Response ---\n{response.json()}\n--- End Submission Response ---\n")

        if data.get("correct") and data.get("url"):
            if data.get("correct") is True:
                print(f"--- Answer correct. Proceeding to next URL: {data.get('url')} ---\n")
                # change the status of last question to True
                questions_data[-1]["status"] = True
                questions_data[-1]["completed_at"] = datetime.now(IST)
                questions_data.append({
                    "url": data.get('url', None),
                    "question_number": 1 + len(questions_data),
                    "status": False,
                    "started_at": datetime.now(IST),
                    "completed_at": None,
                    "answer": None
                })
                payload = await process_task(data.get("url"))
            else:
                last_question = questions_data[-1]
                time_diff = datetime.now(IST) - last_question["started_at"]
                if time_diff.total_seconds() > 60:
                    print(f"--- Time exceeded for {last_question['url']}. Proceeding to next URL: {data.get('url')} ---\n")
                    questions_data[-1]["status"] = False
                    questions_data[-1]["completed_at"] = datetime.now(IST)
                    questions_data.append({
                        "url": data.get('url', None),
                        "question_number": 1 + len(questions_data),
                        "status": False,
                        "started_at": datetime.now(IST),
                        "completed_at": None,
                        "answer": None
                    })
                    payload = await process_task(data.get("url"))
                else:
                    if data.get("reason"):
                        print(f"--- Answer incorrect. Reason: {data.get('reason')} ---\n")
                    payload = await process_task(last_question["url"], data.get("reason"))
        elif data.get("correct") and not data.get("url"):
            if data.get('correct') is False:
                last_question = questions_data[-1]
                time_diff = datetime.now(IST) - last_question["started_at"]
                if time_diff.total_seconds() > 150:
                    print(f"--- Time exceeded for {last_question['url']}. Ending evaluation. ---\n")
                    questions_data[-1]["status"] = False
                    questions_data[-1]["completed_at"] = datetime.now(IST)
                    print(f"--- QUESTIONS SUMMARY ---")
                    for q in questions_data:
                        print(f"Question {q['question_number']}: URL: {q['url']}, Status: {'Correct' if q['status'] else 'Incorrect'}, Started at: {q['started_at'].strftime('%Y-%m-%d %H:%M:%S')}, Completed at: {q['completed_at'].strftime('%Y-%m-%d %H:%M:%S') if q['completed_at'] else 'N/A'}")  #this code is orginally by Devamm007 (24f2000828@ds.study.iitm.ac.in)
                    print("=== END EVALUATION ===\n")
                    break
                else:
                    print(f"--- Answer incorrect. Reason: {data.get('reason')} ---\n")
                    questions_data[-1]["status"] = False
                    payload = await process_task(last_question["url"], data.get("reason"))
            else:
                print(f"--- Answer correct. Proceeding to next URL: {data.get('url')} ---\n")
                # change the status of last question to True
                questions_data[-1]["status"] = True
                questions_data[-1]["completed_at"] = datetime.now(IST)
                # log all questions data in logging format
                print(f"--- QUESTIONS SUMMARY ---")
                for q in questions_data:
                    print(f"Question {q['question_number']}: URL: {q['url']}, Status: {'Correct' if q['status'] else 'Incorrect'}, Started at: {q['started_at'].strftime('%Y-%m-%d %H:%M:%S')}, Completed at: {q['completed_at'].strftime('%Y-%m-%d %H:%M:%S') if q['completed_at'] else 'N/A'}")
                print("=== END EVALUATION ===\n")
                break
            
        iteration_count += 1
        await asyncio.sleep(1)  # brief pause to avoid overwhelming the server
        if iteration_count >= max_iterations:
            print(f"--- QUESTIONS SUMMARY ---")
            for q in questions_data:
                print(f"Question {q['question_number']}: URL: {q['url']}, Status: {'Correct' if q['status'] else 'Incorrect'}, Started at: {q['started_at'].strftime('%Y-%m-%d %H:%M:%S')}, Completed at: {q['completed_at'].strftime('%Y-%m-%d %H:%M:%S') if q['completed_at'] else 'N/A'}")
            print(f"--- Maximum iterations reached. Ending evaluation. ---\n")
            print("=== END EVALUATION ===\n")
            break
    return None            

@app.post("/handle_task")
async def handle_task(data: dict):
    """
    Endpoint to handle incoming tasks for processing.
    Validates the input and initiates background processing.
    """

    # validate secret, email and JSON validity
    if data.get('secret') and data.get('email') and data.get('url'):
        # validate secret
        if not validate_secret(data.get('secret', '')):
            print("--- Invalid secret detected. ---\n")
            raise HTTPException(status_code=403, detail="Invalid secret")
        
        # validate email
        if not validate_email(data.get('email', '')):
            print("--- Invalid email detected. ---\n")
            raise HTTPException(status_code=400, detail="Invalid email")
        
        # validate url
        url_validation_result = validate_url(data.get("url", ""))
        if type(url_validation_result) != bool:
            print(f"--- URL validation failed: {url_validation_result} ---\n")
            raise HTTPException(status_code=400, detail=url_validation_result)
        
        asyncio.create_task(worker(data.get("url")))
    else:
        print("--- Invalid JSON structure detected. ---\n")
        raise HTTPException(status_code=400, detail="Invalid JSON structure")

    return {"status": "Secret validated. Task is being processed in the background."}

if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=8000)